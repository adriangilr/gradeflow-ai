from __future__ import annotations

import csv
import json
import io
import os
import copy
import re
import shutil
import zipfile
from datetime import datetime, timedelta, timezone
from typing import Any

from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseDownload

from src.auth import get_credentials
from src.config import ensure_directories, get_settings

from src.utils.naming import construir_nombre_portfolio
from src.config import NAMING_MODE, MAX_FOLDER_NAME_LEN


try:
    from PyPDF2 import PdfReader  # type: ignore
except Exception:
    PdfReader = None

try:
    from docx import Document  # type: ignore
except Exception:
    Document = None

try:
    from pptx import Presentation  # type: ignore
except Exception:
    Presentation = None

# ==========================================================
# Configuración general
# ==========================================================

DIAS_RECIENTES = 30
AUTOGRADING_CONFIG_FILENAME = "autograding_rules.json"
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tiff", ".tif"}

DEFAULT_AUTOGRADING_CONFIG = {
    "weights": {
        "delivery_valid": 40,
        "evidence_file_or_text": 20,
        "readable_content": 20,
        "minimum_sufficiency": 20,
    },
    "keywords": {
        "enabled": True,
        "list": ["control", "sistema", "resumen"],
        "minimum_matches": 1,
        "required_for_delivery_valid": False,
    },
    "minimum_sufficiency": {
        "min_words_partial": 10,
        "min_words_full": 50,
        "min_chars_partial": 80,
        "min_chars_full": 300,
        "partial_score": 10,
        "full_score": 20,
    },
    "late_policy": {
        "enabled": True,
        "minor_days_threshold": 5,
        "minor_penalty": 5,
        "major_penalty": 10,
        "fallback_penalty_when_late_without_due_date": 5,
    },
}

# Mapeo simple de MIME -> extensión para corregir nombres raros
MIME_TO_EXT = {
    "application/pdf": ".pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
    "application/zip": ".zip",
    "text/plain": ".txt",
    "text/csv": ".csv",
    "application/json": ".json",
    "image/png": ".png",
    "image/jpeg": ".jpg",
    "image/jpg": ".jpg",
    "image/webp": ".webp",
    "image/gif": ".gif",
    "image/tiff": ".tiff",
}


# ==========================================================
# Utilidades generales
# ==========================================================

def limpiar_nombre_archivo(texto: str) -> str:
    """
    Limpia nombres de archivos o carpetas.
    Evita 'sin_nombre' usando fallback más inteligente.
    """
    if not texto or not texto.strip():
        from datetime import datetime
        return f"archivo_{datetime.now().strftime('%Y%m%d_%H%M%S')}"

    texto = texto.strip()
    reemplazos = {
        "/": "-",
        "\\": "-",
        ":": "-",
        "*": "-",
        "?": "",
        '"': "",
        "<": "(",
        ">": ")",
        "|": "-",
    }

    for viejo, nuevo in reemplazos.items():
        texto = texto.replace(viejo, nuevo)

    import re
    texto = re.sub(r"\s+", " ", texto).strip()

    if not texto:
        from datetime import datetime
        return f"archivo_{datetime.now().strftime('%Y%m%d_%H%M%S')}"

    return texto


def slugify_nombre(texto: str) -> str:
    """
    Convierte texto a formato carpeta amigable:
    minusculas, guiones y sin caracteres raros.
    """
    texto = limpiar_nombre_archivo(texto).lower()
    texto = texto.replace("_", " ")
    texto = re.sub(r"[^a-z0-9]+", "-", texto)
    texto = re.sub(r"-+", "-", texto).strip("-")
    return texto or "sin-nombre"


def construir_slug_curso(texto: str, course_id: str) -> str:
    """
    Genera el nombre visible de la carpeta raíz del curso.
    Ejemplo:
    'Seminario Ing. de SW' + id -> 'Seminario Ing. de SW_840182924161'
    """
    nombre_curso = limpiar_nombre_archivo(texto) or "curso"
    return f"{nombre_curso}_{course_id}"


def construir_slug_actividad(texto: str, actividad_id: str) -> str:
    """
    Genera el nombre visible de la carpeta de la actividad.
    Ejemplo:
    'P01 - Timing plan' + id -> 'P01 - Timing plan_840182924183'
    """
    nombre_actividad = limpiar_nombre_archivo(texto) or "actividad"
    return f"{nombre_actividad}_{actividad_id}"


def obtener_timestamp_carpeta() -> str:
    """
    Timestamp corto y estable para nombres de carpeta.
    Formato: YYYYMMDD_HHMM
    """
    return datetime.now().strftime("%Y%m%d_%H%M")

def asegurar_directorio(path: str) -> None:
    """
    Crea el directorio si no existe.
    """
    os.makedirs(path, exist_ok=True)


def preparar_directorio_salida(path: str, limpiar_si_existe: bool = False) -> str:
    """
    Prepara un directorio de salida controlado.
    Si limpiar_si_existe=True, elimina por completo el contenido previo.
    """
    path = os.path.normpath(path)

    if limpiar_si_existe and os.path.exists(path):
        shutil.rmtree(path)

    asegurar_directorio(path)
    return path


def asegurar_extension(nombre: str, mime_type: str) -> str:
    """
    Garantiza que el nombre tenga una extensión coherente con el MIME.
    Esto evita casos donde Classroom/Drive entrega títulos raros como .pod
    o nombres sin extensión.
    """
    base, ext = os.path.splitext(nombre)
    ext_actual = ext.lower()
    ext_correcta = MIME_TO_EXT.get((mime_type or "").lower())

    if ext_correcta is None:
        return nombre

    if not ext_actual:
        return f"{base}{ext_correcta}"

    if ext_actual != ext_correcta:
        return f"{base}{ext_correcta}"

    return nombre


def construir_nombre_carpeta_entrega(
    submission: dict,
    perfil: dict,
) -> str:

    nombre = perfil.get("nombre", "")
    apellido = perfil.get("apellido", "")
    user_id = str(submission.get("userId", ""))

    return construir_nombre_portfolio(
        nombre=nombre,
        apellido=apellido,
        user_id=user_id,
        modo=NAMING_MODE,
        max_len=MAX_FOLDER_NAME_LEN,
    )


def seleccionar_opcion(lista: list[dict[str, Any]], tipo: str) -> dict[str, Any]:
    """
    Menú interactivo genérico para terminal.
    """
    print(f"\nSelecciona {tipo}:\n")

    for i, item in enumerate(lista, start=1):
        print(f"{i}. {item['display_name']}")

    while True:
        entrada = input("\nIngresa número: ").strip()

        try:
            indice = int(entrada) - 1
            if 0 <= indice < len(lista):
                return lista[indice]
        except ValueError:
            pass

        print("❌ Opción inválida. Intenta de nuevo.")


def parse_google_datetime(value: str | None) -> datetime | None:
    """
    Convierte timestamps de Google tipo:
    2026-04-09T18:25:43.123Z
    """
    if not value:
        return None

    try:
        if value.endswith("Z"):
            value = value.replace("Z", "+00:00")
        return datetime.fromisoformat(value)
    except ValueError:
        return None


def ahora_utc() -> datetime:
    """
    Regresa la fecha actual en UTC.
    """
    return datetime.now(timezone.utc)


def cargar_config_autograding() -> dict[str, Any]:
    """
    Carga reglas de autograding desde JSON.
    Si no existe o falla, usa la configuración por defecto.
    """
    posibles_rutas = [
        os.path.join(os.path.dirname(__file__), AUTOGRADING_CONFIG_FILENAME),
        os.path.join(os.getcwd(), AUTOGRADING_CONFIG_FILENAME),
    ]

    for ruta in posibles_rutas:
        if os.path.exists(ruta):
            try:
                with open(ruta, "r", encoding="utf-8") as f:
                    data = json.load(f)

                if isinstance(data, dict):
                    return combinar_config(DEFAULT_AUTOGRADING_CONFIG, data)
            except Exception as err:
                print(f"⚠️ No se pudo leer la configuración de autograding en '{ruta}': {err}")

    return copy.deepcopy(DEFAULT_AUTOGRADING_CONFIG)


def combinar_config(base: dict[str, Any], extra: dict[str, Any]) -> dict[str, Any]:
    """Mezcla profunda simple de diccionarios."""
    resultado: dict[str, Any] = {}

    for key, value in base.items():
        if isinstance(value, dict):
            extra_value = extra.get(key, {})
            if isinstance(extra_value, dict):
                resultado[key] = combinar_config(value, extra_value)
            else:
                resultado[key] = copy.deepcopy(value)
        else:
            resultado[key] = extra.get(key, value)

    for key, value in extra.items():
        if key not in resultado:
            resultado[key] = value

    return resultado


# ==========================================================
# Nombres visibles para menús
# ==========================================================

def obtener_nombre_curso_visible(course: dict[str, Any]) -> str:
    """
    Construye un nombre visible único para evitar confusiones
    cuando existen dos cursos con el mismo nombre.
    """
    name = course.get("name", "Curso sin nombre")
    section = course.get("section", "").strip()
    course_id = course.get("id", "").strip()
    room = course.get("room", "").strip()

    extras = []
    if section:
        extras.append(section)
    if room:
        extras.append(f"Aula: {room}")
    if course_id:
        extras.append(f"id={course_id}")

    return f"{name} | {' | '.join(extras)}" if extras else name


def obtener_nombre_actividad_visible(coursework: dict[str, Any]) -> str:
    """
    Construye un nombre más legible para la actividad.
    """
    title = coursework.get("title", "Actividad sin título")
    coursework_id = coursework.get("id", "")
    work_type = coursework.get("workType", "")
    max_points = coursework.get("maxPoints")

    extras = []
    if work_type:
        extras.append(work_type)
    if max_points is not None:
        extras.append(f"{max_points} pts")
    if coursework_id:
        extras.append(f"id={coursework_id}")

    return f"{title} | {' | '.join(extras)}" if extras else title


# ==========================================================
# Menús
# ==========================================================

def seleccionar_alcance_descarga() -> str:
    """
    Permite elegir si se descargará una actividad
    o todas las actividades del curso.
    """
    opciones = [
        {
            "id": "single_coursework",
            "display_name": "Descargar una sola actividad",
        },
        {
            "id": "all_courseworks",
            "display_name": "Descargar todas las actividades del curso",
        },
    ]
    return seleccionar_opcion(opciones, "el alcance de la descarga")["id"]


def seleccionar_modo_descarga() -> str:
    """
    Define qué entregas se descargarán.
    """
    opciones = [
        {"id": "all", "display_name": "Incluir todos los alumnos (entregados y no entregados)"},
        {
            "id": "late_ungraded",
            "display_name": "Bajar solo tardías y no evaluadas",
        },
    ]
    return seleccionar_opcion(opciones, "un modo de descarga")["id"]


def seleccionar_formato_salida() -> str:
    """
    Define si solo se guarda en carpeta o también en zip.
    """
    opciones = [
        {"id": "folder_only", "display_name": "Guardar solo en carpeta"},
        {"id": "zip_and_folder", "display_name": "Guardar en carpeta y generar .zip"},
    ]
    return seleccionar_opcion(opciones, "un formato de salida")["id"]


def seleccionar_filtro_actividades() -> str:
    """
    Permite aplicar un filtro a las actividades del curso
    antes de elegir una o antes de procesarlas todas.
    """
    opciones = [
        {"id": "all", "display_name": "Todas las actividades"},
        {"id": "with_submissions", "display_name": "Solo actividades con entregas"},
    ]
    return seleccionar_opcion(opciones, "un filtro de actividades")["id"]


def describir_modo_descarga(modo_descarga: str) -> str:
    descripciones = {
        "all": "todas las entregas activas",
        "resubmitted": "solo reentregadas y pendientes de reevaluación",
        "ungraded": "solo no evaluadas",
        "late": "solo tardías",
        "resubmitted_ungraded": "solo reentregadas y no evaluadas",
        "late_ungraded": "solo tardías y no evaluadas",
    }
    return descripciones.get(modo_descarga, "filtro desconocido")


# ==========================================================
# Descarga de Drive
# ==========================================================

def obtener_metadata_archivo_drive(drive_service, file_id: str) -> dict[str, str]:
    """
    Lee metadata real del archivo en Drive para usar el nombre correcto
    y el mimeType real, en lugar de confiar ciegamente en 'title'.
    """
    try:
        meta = (
            drive_service.files()
            .get(fileId=file_id, fields="id,name,mimeType,fileExtension")
            .execute()
        )
        return {
            "name": meta.get("name", ""),
            "mimeType": meta.get("mimeType", ""),
            "fileExtension": meta.get("fileExtension", ""),
        }
    except HttpError:
        return {
            "name": "",
            "mimeType": "",
            "fileExtension": "",
        }


def download_file(
    drive_service,
    file_id: str,
    file_name: str,
    folder: str,
    mime_type: str = "",
) -> str | None:
    """
    Descarga un archivo de Drive al folder indicado.
    Usa metadata real de Drive para corregir extensiones raras o faltantes.
    """
    asegurar_directorio(folder)

    try:
        meta = obtener_metadata_archivo_drive(drive_service, file_id)

        # Fuente de verdad:
        # 1) name real de Drive
        # 2) si no viene, usar file_name recibido
        real_name = meta.get("name") or file_name or "archivo"
        real_mime = meta.get("mimeType") or mime_type or ""

        safe_name = limpiar_nombre_archivo(real_name)
        safe_name = asegurar_extension(safe_name, real_mime)

        file_path = os.path.join(folder, safe_name)

        request = drive_service.files().get_media(fileId=file_id)
        with io.FileIO(file_path, "wb") as fh:
            downloader = MediaIoBaseDownload(fh, request)

            done = False
            while not done:
                _, done = downloader.next_chunk()

        print(f"      ✅ Descargado: {safe_name}")
        return file_path

    except HttpError as err:
        print(f"      ❌ Error al descargar '{file_name}': {err}")
        return None


# ==========================================================
# Lectura de Classroom
# ==========================================================

def obtener_todos_los_cursos(classroom_service) -> list[dict[str, Any]]:
    """
    Recupera todos los cursos paginando.
    Si existen cursos activos, prioriza esos.
    """
    courses: list[dict[str, Any]] = []
    page_token = None

    while True:
        response = (
            classroom_service.courses()
            .list(pageSize=100, pageToken=page_token)
            .execute()
        )

        courses.extend(response.get("courses", []))
        page_token = response.get("nextPageToken")

        if not page_token:
            break

    cursos_activos = [c for c in courses if c.get("courseState", "ACTIVE") == "ACTIVE"]
    cursos_finales = cursos_activos if cursos_activos else courses

    for course in cursos_finales:
        course["display_name"] = obtener_nombre_curso_visible(course)

    cursos_finales.sort(key=lambda x: x.get("display_name", "").lower())
    return cursos_finales


def obtener_todas_las_actividades(classroom_service, course_id: str) -> list[dict[str, Any]]:
    """
    Recupera todas las actividades de un curso.
    """
    courseworks: list[dict[str, Any]] = []
    page_token = None

    while True:
        response = (
            classroom_service.courses()
            .courseWork()
            .list(courseId=course_id, pageSize=100, pageToken=page_token)
            .execute()
        )

        courseworks.extend(response.get("courseWork", []))
        page_token = response.get("nextPageToken")

        if not page_token:
            break

    for coursework in courseworks:
        coursework["display_name"] = obtener_nombre_actividad_visible(coursework)

    courseworks.sort(key=lambda x: x.get("display_name", "").lower())
    return courseworks


def obtener_todas_las_entregas(
    classroom_service,
    course_id: str,
    coursework_id: str,
) -> list[dict[str, Any]]:
    """
    Recupera todas las entregas de una actividad.
    """
    submissions: list[dict[str, Any]] = []
    page_token = None

    while True:
        response = (
            classroom_service.courses()
            .courseWork()
            .studentSubmissions()
            .list(
                courseId=course_id,
                courseWorkId=coursework_id,
                pageSize=100,
                pageToken=page_token,
            )
            .execute()
        )

        submissions.extend(response.get("studentSubmissions", []))
        page_token = response.get("nextPageToken")

        if not page_token:
            break

    return submissions


# ==========================================================
# Filtros de entregas
# ==========================================================

def ya_fue_devuelta_antes(submission: dict[str, Any]) -> bool:
    """
    Detecta si la entrega ya había sido devuelta antes.
    """
    for event in submission.get("submissionHistory", []):
        state_history = event.get("stateHistory", {})
        if state_history.get("state") == "RETURNED":
            return True
    return False


def es_reentregada(submission: dict[str, Any]) -> bool:
    return (
        submission.get("state") == "TURNED_IN"
        and ya_fue_devuelta_antes(submission)
    )


def es_no_evaluada(submission: dict[str, Any]) -> bool:
    return (
        submission.get("state") == "TURNED_IN"
        and submission.get("assignedGrade") is None
    )


def es_tardia(submission: dict[str, Any]) -> bool:
    return (
        submission.get("state") == "TURNED_IN"
        and submission.get("late", False) is True
    )


def estado_legible_entrega(submission: dict[str, Any]) -> str:
    """
    Traduce el estado técnico de Classroom a algo más entendible.
    """
    state = (submission.get("state") or "").upper()

    mapa = {
        "TURNED_IN": "entregado",
        "CREATED": "asignado_sin_entregar",
        "RETURNED": "devuelto",
        "RECLAIMED_BY_STUDENT": "reclamado_por_alumno",
    }
    return mapa.get(state, state.lower() or "desconocido")


def se_puede_descargar_entrega(submission: dict[str, Any]) -> bool:
    """
    Solo tiene sentido intentar descargar adjuntos cuando la entrega fue enviada.
    """
    return submission.get("state") == "TURNED_IN"


def filtrar_entregas(submissions: list[dict[str, Any]], modo_descarga: str) -> list[dict[str, Any]]:
    """
    Aplica el filtro elegido a la lista de entregas.
    """
    if modo_descarga == "all":
        # Incluye a todos los alumnos de la actividad:
        # entregados, asignados sin entregar, devueltos, etc.
        return submissions

    if modo_descarga == "resubmitted":
        return [s for s in submissions if es_reentregada(s)]

    if modo_descarga == "ungraded":
        return [s for s in submissions if es_no_evaluada(s)]

    if modo_descarga == "late":
        return [s for s in submissions if es_tardia(s)]

    if modo_descarga == "resubmitted_ungraded":
        return [s for s in submissions if es_reentregada(s) and es_no_evaluada(s)]

    if modo_descarga == "late_ungraded":
        return [s for s in submissions if es_tardia(s) and es_no_evaluada(s)]

    return []


# ==========================================================
# Filtros de actividades
# ==========================================================

def actividad_publicada(coursework: dict[str, Any]) -> bool:
    """
    Considera publicada cuando state es PUBLISHED.
    Si el campo no viene, asumimos True para no esconder actividades válidas.
    """
    state = coursework.get("state")
    if state is None:
        return True
    return state == "PUBLISHED"


def actividad_reciente(coursework: dict[str, Any], dias: int = DIAS_RECIENTES) -> bool:
    """
    Considera reciente una actividad creada o actualizada dentro
    de los últimos N días.
    """
    limite = ahora_utc() - timedelta(days=dias)

    creation_time = parse_google_datetime(coursework.get("creationTime"))
    update_time = parse_google_datetime(coursework.get("updateTime"))
    due_date = None

    due = coursework.get("dueDate")
    if isinstance(due, dict):
        try:
            year = due.get("year")
            month = due.get("month")
            day = due.get("day")
            if year and month and day:
                due_date = datetime(year, month, day, tzinfo=timezone.utc)
        except ValueError:
            due_date = None

    fechas_validas = [f for f in [creation_time, update_time, due_date] if f is not None]
    if not fechas_validas:
        return False

    return any(f >= limite for f in fechas_validas)


def filtrar_actividades(
    classroom_service,
    course_id: str,
    courseworks: list[dict[str, Any]],
    filtro: str,
) -> list[dict[str, Any]]:
    """
    Aplica filtro a las actividades.
    """
    if filtro == "all":
        return courseworks

    if filtro == "published":
        return [cw for cw in courseworks if actividad_publicada(cw)]

    if filtro == "recent":
        return [cw for cw in courseworks if actividad_reciente(cw)]

    if filtro == "with_submissions":
        filtradas = []
        for cw in courseworks:
            try:
                submissions = obtener_todas_las_entregas(
                    classroom_service=classroom_service,
                    course_id=course_id,
                    coursework_id=cw["id"],
                )
                if submissions:
                    filtradas.append(cw)
            except HttpError as err:
                print(
                    f"⚠️ No se pudieron revisar entregas para actividad "
                    f"'{cw.get('title', cw.get('id', 'sin_titulo'))}': {err}"
                )
        return filtradas

    return courseworks


# ==========================================================
# Perfil del alumno
# ==========================================================

def extraer_datos_usuario_desde_historial(submission: dict[str, Any]) -> dict[str, str]:
    """
    Intenta recuperar nombre/correo desde submissionHistory.
    Esto sirve como plan B cuando no hay scopes suficientes
    para consultar userProfiles.
    """
    history = submission.get("submissionHistory", [])

    for event in history:
        actor = event.get("actorUser", {})
        if not actor:
            continue

        name = actor.get("name", {}) or {}
        given_name = name.get("givenName", "") or ""
        family_name = name.get("familyName", "") or ""
        full_name = name.get("fullName", "") or ""

        email = actor.get("emailAddress", "") or ""

        if not given_name and full_name:
            partes = full_name.split()
            if partes:
                given_name = partes[0]
                if len(partes) > 1:
                    family_name = " ".join(partes[1:])

        if given_name or family_name or email:
            return {
                "correo": email,
                "nombre": given_name,
                "apellido": family_name,
            }

    return {
        "correo": "",
        "nombre": "",
        "apellido": "",
    }


def obtener_perfil_usuario(
    classroom_service,
    user_id: str,
    profile_scope_disponible: bool,
) -> dict[str, str]:
    """
    Recupera correo, nombre y apellido del alumno.

    Estrategia:
    1. Si hay scope disponible, intenta userProfiles
    2. Si no hay scope, el caller puede usar fallback desde historial
    """
    if not profile_scope_disponible:
        return {
            "correo": "",
            "nombre": "",
            "apellido": "",
        }

    try:
        profile = classroom_service.userProfiles().get(userId=user_id).execute()

        email = profile.get("emailAddress", "") or ""

        name = profile.get("name", {}) or {}
        given_name = name.get("givenName", "") or ""
        family_name = name.get("familyName", "") or ""
        full_name = name.get("fullName", "") or ""

        if not given_name and full_name:
            partes = full_name.split()
            if partes:
                given_name = partes[0]
                if len(partes) > 1:
                    family_name = " ".join(partes[1:])

        return {
            "correo": email,
            "nombre": given_name,
            "apellido": family_name,
        }

    except HttpError as err:
        raise err


def detectar_scope_perfil(classroom_service) -> bool:
    """
    Prueba una sola vez si el token tiene permisos para consultar perfiles.
    Así evitamos un error 403 repetido para cada alumno.
    """
    try:
        classroom_service.userProfiles().get(userId="me").execute()
        return True
    except HttpError as err:
        status = getattr(err, "status_code", None)
        contenido = str(err)

        if status == 403 or "ACCESS_TOKEN_SCOPE_INSUFFICIENT" in contenido:
            return False

        raise err


# ==========================================================
# Adjuntos
# ==========================================================

def obtener_adjuntos(submission: dict[str, Any]) -> list[dict[str, Any]]:
    assignment_submission = submission.get("assignmentSubmission", {})
    return assignment_submission.get("attachments", [])


def tiene_adjuntos(submission: dict[str, Any]) -> bool:
    return len(obtener_adjuntos(submission)) > 0


def obtener_due_date_texto(coursework: dict[str, Any]) -> str:
    """
    Convierte dueDate de Classroom a texto YYYY-MM-DD.
    """
    due_date = coursework.get("dueDate")
    if not isinstance(due_date, dict):
        return ""

    year = due_date.get("year")
    month = due_date.get("month")
    day = due_date.get("day")

    if not (year and month and day):
        return ""

    try:
        return f"{int(year):04d}-{int(month):02d}-{int(day):02d}"
    except (TypeError, ValueError):
        return ""


def obtener_due_time_texto(coursework: dict[str, Any]) -> str:
    """
    Convierte dueTime de Classroom a texto HH:MM:SS.
    """
    due_time = coursework.get("dueTime")
    if not isinstance(due_time, dict):
        return ""

    hours = due_time.get("hours", 0)
    minutes = due_time.get("minutes", 0)
    seconds = due_time.get("seconds", 0)

    try:
        return f"{int(hours):02d}:{int(minutes):02d}:{int(seconds):02d}"
    except (TypeError, ValueError):
        return ""


# ==========================================================
# Lectura y evaluación de contenido
# ==========================================================

def leer_texto_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def leer_texto_pdf(path: str) -> str:
    if PdfReader is None:
        return ""

    try:
        reader = PdfReader(path)
        partes: list[str] = []
        for page in reader.pages:
            partes.append(page.extract_text() or "")
        return "\n".join(partes)
    except Exception:
        return ""


def leer_texto_docx(path: str) -> str:
    if Document is None:
        return ""

    try:
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text)
    except Exception:
        return ""


def leer_texto_zip(path: str, profundidad_max: int = 15) -> str:
    """
    Intenta leer texto útil de archivos simples dentro de un ZIP.
    No revienta si el ZIP trae binarios.
    """
    partes: list[str] = []

    try:
        with zipfile.ZipFile(path, "r") as zf:
            for idx, name in enumerate(zf.namelist()):
                if idx >= profundidad_max:
                    break

                lower = name.lower()
                if lower.endswith((".txt", ".md", ".csv", ".py", ".json", ".log")):
                    try:
                        data = zf.read(name)
                        partes.append(data.decode("utf-8", errors="ignore"))
                    except Exception:
                        continue
    except Exception:
        return ""

    return "\n".join(partes)


def es_archivo_imagen(path: str) -> bool:
    """
    Detecta si el archivo es una imagen común.
    """
    ext = os.path.splitext(path)[1].lower()
    return ext in IMAGE_EXTENSIONS


def contiene_imagenes(rutas: list[str]) -> bool:
    """
    Indica si entre los adjuntos descargados hay al menos una imagen.
    """
    return any(es_archivo_imagen(ruta) for ruta in rutas)


def leer_texto_pptx(path: str) -> str:
    if Presentation is None:
        return ""

    try:
        prs = Presentation(path)
        textos = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    textos.append(shape.text)

        return "\n".join(textos)
    except Exception:
        return ""


def extraer_texto_archivo(path: str) -> str:
    """
    Lee texto de varios tipos de archivo comunes.
    Si no puede, regresa cadena vacía.
    """
    ext = os.path.splitext(path)[1].lower()

    if ext in {".txt", ".md", ".csv", ".py", ".json", ".log"}:
        return leer_texto_txt(path)

    if ext == ".pdf":
        return leer_texto_pdf(path)

    if ext == ".docx":
        return leer_texto_docx(path)

    if ext == ".zip":
        return leer_texto_zip(path)

    if ext == ".pptx":
        return leer_texto_pptx(path)

    if ext in IMAGE_EXTENSIONS:
        return ""

    return ""


def analizar_contenido_texto(
    texto: str,
    autograding_config: dict[str, Any],
) -> dict[str, Any]:
    """
    Analiza el texto extraído con una lógica más simple:
    - detecta si el contenido es legible
    - calcula suficiencia mínima configurable
    - detecta palabras clave opcionales
    """
    texto_limpio = re.sub(r"\s+", " ", texto or "").strip()
    palabras = re.findall(r"\b\w+\b", texto_limpio, flags=re.UNICODE)
    num_palabras = len(palabras)
    num_caracteres = len(texto_limpio)

    weights = autograding_config.get("weights", {})
    suff_cfg = autograding_config.get("minimum_sufficiency", {})
    keywords_cfg = autograding_config.get("keywords", {})

    contenido_legible = bool(texto_limpio)

    min_words_partial = int(suff_cfg.get("min_words_partial", 10))
    min_words_full = int(suff_cfg.get("min_words_full", 50))
    min_chars_partial = int(suff_cfg.get("min_chars_partial", 80))
    min_chars_full = int(suff_cfg.get("min_chars_full", 300))
    partial_score = int(suff_cfg.get("partial_score", 10))
    full_score = int(suff_cfg.get("full_score", weights.get("minimum_sufficiency", 20)))

    if num_palabras >= min_words_full or num_caracteres >= min_chars_full:
        sufficiency_score = full_score
        sufficiency_level = "full"
    elif num_palabras >= min_words_partial or num_caracteres >= min_chars_partial:
        sufficiency_score = partial_score
        sufficiency_level = "partial"
    else:
        sufficiency_score = 0
        sufficiency_level = "low"

    keyword_list = keywords_cfg.get("list", [])
    if not isinstance(keyword_list, list):
        keyword_list = []

    keyword_hits: list[str] = []
    if keywords_cfg.get("enabled", True):
        texto_lower = texto_limpio.lower()
        for keyword in keyword_list:
            kw = str(keyword).strip().lower()
            if kw and kw in texto_lower:
                keyword_hits.append(kw)

    minimum_matches = int(keywords_cfg.get("minimum_matches", 1))
    keywords_ok = len(keyword_hits) >= minimum_matches if keyword_list else True

    return {
        "texto_extraido": texto_limpio,
        "num_palabras": num_palabras,
        "num_caracteres": num_caracteres,
        "contenido_legible": contenido_legible,
        "sufficiency_score": sufficiency_score,
        "sufficiency_level": sufficiency_level,
        "keyword_hits": keyword_hits,
        "keywords_ok": keywords_ok,
    }


def construir_due_datetime(coursework: dict[str, Any]) -> datetime | None:
    due_date = coursework.get("dueDate")
    if not isinstance(due_date, dict):
        return None

    year = due_date.get("year")
    month = due_date.get("month")
    day = due_date.get("day")
    if not (year and month and day):
        return None

    due_time = coursework.get("dueTime") or {}
    hours = due_time.get("hours", 23)
    minutes = due_time.get("minutes", 59)
    seconds = due_time.get("seconds", 59)

    try:
        return datetime(
            int(year),
            int(month),
            int(day),
            int(hours),
            int(minutes),
            int(seconds),
            tzinfo=timezone.utc,
        )
    except (TypeError, ValueError):
        return None


def obtener_timestamp_entrega(submission: dict[str, Any]) -> datetime | None:
    candidatos = [
        submission.get("updateTime"),
        submission.get("submissionTime"),
        submission.get("turnInTime"),
        submission.get("creationTime"),
    ]
    for valor in candidatos:
        dt = parse_google_datetime(valor)
        if dt is not None:
            return dt
    return None


def calcular_penalizacion_tardanza(
    submission: dict[str, Any],
    coursework: dict[str, Any],
    autograding_config: dict[str, Any],
) -> tuple[int, int]:
    late_policy = autograding_config.get("late_policy", {})
    if not late_policy.get("enabled", True):
        return 0, 0

    if not bool(submission.get("late", False)):
        return 0, 0

    due_dt = construir_due_datetime(coursework)
    entrega_dt = obtener_timestamp_entrega(submission)

    if due_dt is None or entrega_dt is None:
        penalty = int(late_policy.get("fallback_penalty_when_late_without_due_date", 5))
        return penalty, 0

    delta = entrega_dt - due_dt
    total_dias = max(0, int((delta.total_seconds() + 86399) // 86400))

    threshold = int(late_policy.get("minor_days_threshold", 5))
    minor_penalty = int(late_policy.get("minor_penalty", 5))
    major_penalty = int(late_policy.get("major_penalty", 10))

    penalty = minor_penalty if total_dias <= threshold else major_penalty
    return penalty, total_dias


def construir_feedback(
    late: bool,
    has_attachment: bool,
    archivos_leidos: int,
    num_palabras: int,
    penalty_late: int,
    auto_grade: int,
    manual_review: bool,
    contenido_legible: bool,
    sufficiency_level: str,
    keyword_hits: list[str],
) -> str:
    """
    Genera feedback más limpio y profesional.
    Evita mensajes basura cuando sí existe evidencia, pero no fue interpretable.
    """
    mensajes: list[str] = []

    # 1) Evidencia
    if has_attachment:
        mensajes.append("Se recibió evidencia de entrega.")
    else:
        mensajes.append("No se detectó evidencia adjunta ni texto interpretable.")

    # 2) Tardanza
    if late and penalty_late > 0:
        mensajes.append(f"Se aplicó penalización por tardanza (-{penalty_late}).")
    else:
        mensajes.append("Sin penalización por tardanza.")

    # 3) Estado de lectura
    if manual_review:
        mensajes.append(
            "La evidencia requiere revisión manual porque incluye imágenes o formatos no interpretables automáticamente."
        )
    elif archivos_leidos > 0 and contenido_legible:
        mensajes.append(
            f"Se detectó contenido legible automáticamente ({num_palabras} palabras aprox.)."
        )
    elif has_attachment:
        mensajes.append(
            "Se recibió evidencia, pero no fue posible interpretarla automáticamente con las librerías actuales."
        )
    else:
        mensajes.append(
            "No fue posible recuperar contenido legible para evaluación automática."
        )

    # 4) Keywords: solo se informan si sí existen
    if keyword_hits:
        mensajes.append(f"Palabras clave detectadas: {', '.join(keyword_hits)}.")

    # 5) Suficiencia
    if sufficiency_level == "full":
        mensajes.append("El contenido cumple suficiencia completa según los umbrales configurados.")
    elif sufficiency_level == "partial":
        mensajes.append("El contenido cumple suficiencia mínima parcial; conviene revisión rápida.")
    else:
        mensajes.append("El contenido detectado es breve o insuficiente según los umbrales configurados.")

    # 6) Cierre
    mensajes.append(f"Calificación automática sugerida: {auto_grade}/100.")
    return " ".join(mensajes)


def construir_feedback_corto(
    late: bool,
    manual_review: bool,
    contenido_legible: bool,
    sufficiency_level: str,
) -> str:
    """
    Genera feedback compacto tipo etiquetas para CSV y dashboards.
    Ejemplos:
    - suficiente
    - parcial
    - insuficiente
    - tardía | manual_review | insuficiente
    """
    tags: list[str] = []

    if late:
        tags.append("tardía")

    if manual_review:
        tags.append("manual_review")
    elif not contenido_legible:
        tags.append("no_legible")

    if sufficiency_level == "full":
        tags.append("suficiente")
    elif sufficiency_level == "partial":
        tags.append("parcial")
    else:
        tags.append("insuficiente")

    if not tags:
        return "valida"

    return " | ".join(tags)


def evaluar_entrega_automatica(
    submission: dict[str, Any],
    rutas_descargadas: list[str],
    coursework: dict[str, Any],
    autograding_config: dict[str, Any],
) -> dict[str, Any]:
    """
    Evalúa una entrega con estrategia simple y configurable.
    """
    weights = autograding_config.get("weights", {})
    keywords_cfg = autograding_config.get("keywords", {})

    late = bool(submission.get("late", False))
    entrego = submission.get("state") == "TURNED_IN"
    has_attachment = len(rutas_descargadas) > 0 or tiene_adjuntos(submission)

    texto_total: list[str] = []
    archivos_leidos = 0
    manual_review = contiene_imagenes(rutas_descargadas)

    for ruta in rutas_descargadas:
        texto = extraer_texto_archivo(ruta)
        if texto.strip():
            texto_total.append(texto)
            archivos_leidos += 1

    texto_unido = "\n".join(texto_total)
    analisis = analizar_contenido_texto(texto_unido, autograding_config=autograding_config)

    penalty_late, days_late = calcular_penalizacion_tardanza(
        submission=submission,
        coursework=coursework,
        autograding_config=autograding_config,
    )

    keyword_required_for_delivery = bool(
        keywords_cfg.get("required_for_delivery_valid", False)
    )

    delivery_valid_score = 0
    if entrego and (not keyword_required_for_delivery or analisis["keywords_ok"]):
        delivery_valid_score = int(weights.get("delivery_valid", 40))

    evidence_score = (
        int(weights.get("evidence_file_or_text", 20))
        if (has_attachment or analisis["contenido_legible"])
        else 0
    )
    readable_content_score = (
        int(weights.get("readable_content", 20))
        if analisis["contenido_legible"]
        else 0
    )
    sufficiency_score = int(analisis["sufficiency_score"])

    if not entrego:
        auto_grade = 0
        manual_review = False
    else:
        auto_grade = (
            delivery_valid_score
            + evidence_score
            + readable_content_score
            + sufficiency_score
            - penalty_late
        )
        auto_grade = max(0, min(100, auto_grade))

    feedback = construir_feedback_corto(
        late=late,
        manual_review=manual_review,
        contenido_legible=bool(analisis["contenido_legible"]),
        sufficiency_level=str(analisis["sufficiency_level"]),
    )

    return {
        "auto_grade": auto_grade,
        "feedback": feedback,
        "penalty_late": penalty_late,
        "days_late": days_late,
        "has_attachment": str(has_attachment).lower(),
        "delivery_valid_score": delivery_valid_score,
        "evidence_score": evidence_score,
        "readable_content_score": readable_content_score,
        "content_score": sufficiency_score,
        "minimum_sufficiency_score": sufficiency_score,
        "files_read_for_content": archivos_leidos,
        "detected_words": int(analisis["num_palabras"]),
        "detected_characters": int(analisis["num_caracteres"]),
        "keyword_hits": ", ".join(analisis["keyword_hits"]),
        "manual_review": str(manual_review).lower(),
        "readable_content": str(bool(analisis["contenido_legible"])).lower(),
    }


def imprimir_resumen_entrega(submission: dict[str, Any]) -> None:
    """
    Imprime información útil para la revisión en terminal.
    """
    print(f"Entrega ID: {submission.get('id', 'Sin ID')}")
    print(f"  userId: {submission.get('userId', 'Sin userId')}")
    print(f"  estado: {submission.get('state', 'Sin estado')}")
    print(f"  late: {submission.get('late', False)}")
    print(f"  assignedGrade: {submission.get('assignedGrade')}")
    print(f"  draftGrade: {submission.get('draftGrade')}")
    print(f"  reentregada: {es_reentregada(submission)}")
    print(f"  no evaluada: {es_no_evaluada(submission)}")
    print(f"  attached: {tiene_adjuntos(submission)}")


def descargar_adjuntos_entrega(
    submission: dict[str, Any],
    drive_service,
    carpeta_entrega: str,
) -> list[str]:
    """
    Descarga adjuntos de una entrega.
    Regresa la lista de archivos reales descargados.
    """
    attachments = obtener_adjuntos(submission)
    rutas_descargadas: list[str] = []

    if not attachments:
        print("  adjuntos: ninguno")
        return rutas_descargadas

    print("  adjuntos:")

    for att in attachments:
        if "driveFile" in att:
            drive_file = att.get("driveFile", {})
            drive_meta = drive_file.get("driveFile") or drive_file

            file_id = drive_meta.get("id")
            title = drive_meta.get("title", "archivo")
            mime_type = drive_meta.get("mimeType", "")

            print(f"    - DriveFile: {title} | id={file_id}")

            if file_id:
                ruta = download_file(
                    drive_service=drive_service,
                    file_id=file_id,
                    file_name=title,
                    folder=carpeta_entrega,
                    mime_type=mime_type,
                )
                if ruta:
                    rutas_descargadas.append(ruta)

        elif "link" in att:
            link = att["link"]
            print(
                f"    - Link: {link.get('title', 'Sin título')} "
                f"| url={link.get('url', 'Sin URL')}"
            )

        elif "form" in att:
            form = att["form"]
            print(
                f"    - Form: {form.get('title', 'Sin título')} "
                f"| url={form.get('formUrl', 'Sin URL')}"
            )

        elif "youTubeVideo" in att:
            video = att["youTubeVideo"]
            print(
                f"    - YouTube: {video.get('title', 'Sin título')} "
                f"| url={video.get('alternateLink', 'Sin URL')}"
            )

        else:
            print("    - Tipo de adjunto no manejado directamente.")

    return rutas_descargadas


# ==========================================================
# CSV y ZIP
# ==========================================================

def escribir_csv_resumen(csv_path: str, filas: list[dict[str, str]]) -> None:
    """
    Genera CSV general con información de curso, actividad, alumno
    y evaluación automática.
    """
    asegurar_directorio(os.path.dirname(csv_path))

    with open(csv_path, "w", newline="", encoding="utf-8") as csvfile:
        writer = csv.DictWriter(
            csvfile,
            fieldnames=[
                "curso",
                "actividad",
                "due_date",
                "due_time",
                "correo",
                "nombre",
                "apellido",
                "estado_entrega",
                "late",
                "assigned_grade",
                "draft_grade",
                "attached",
                "has_attachment",
                "manual_review",
                "readable_content",
                "days_late",
                "penalty_late",
                "delivery_valid_score",
                "evidence_score",
                "readable_content_score",
                "minimum_sufficiency_score",
                "keyword_hits",
                "content_score",
                "auto_grade",
                "feedback",
            ],
        )
        writer.writeheader()
        writer.writerows(filas)

    print(f"\n✅ CSV generado: {csv_path}")


def comprimir_carpeta_a_zip(carpeta_origen: str, zip_sin_extension: str) -> str:
    """
    Comprime toda la carpeta en un zip.
    """
    zip_path = shutil.make_archive(zip_sin_extension, "zip", carpeta_origen)
    print(f"✅ ZIP generado: {zip_path}")
    return zip_path


# ==========================================================
# Procesamiento de una actividad
# ==========================================================

def procesar_actividad(
    classroom_service,
    drive_service,
    course: dict[str, Any],
    coursework: dict[str, Any],
    modo_descarga: str,
    carpeta_base: str,
    perfiles_cache: dict[str, dict[str, str]],
    filas_csv: list[dict[str, str]],
    estadisticas: dict[str, int],
    profile_scope_disponible: bool,
) -> None:
    """
    Procesa una actividad completa:
    - recupera entregas
    - aplica filtro
    - descarga adjuntos
    - agrega filas al CSV
    - actualiza estadísticas
    """
    course_id = course["id"]
    course_name = course.get("name", f"curso_{course_id}")

    coursework_id = coursework["id"]
    coursework_title = coursework.get("title", f"actividad_{coursework_id}")
    coursework_display = coursework.get("display_name", coursework_title)

    print("\n" + "=" * 90)
    print(f"Procesando actividad: {coursework_display}")
    print("=" * 90)

    submissions = obtener_todas_las_entregas(
        classroom_service=classroom_service,
        course_id=course_id,
        coursework_id=coursework_id,
    )

    estadisticas["actividades_procesadas"] += 1
    estadisticas["entregas_totales"] += len(submissions)

    if not submissions:
        print("No se encontraron entregas para esta actividad.")
        return

    entregas_filtradas = filtrar_entregas(submissions, modo_descarga)
    estadisticas["entregas_filtradas"] += len(entregas_filtradas)

    print(f"Entregas totales encontradas: {len(submissions)}")
    print(f"Entregas a procesar con este filtro: {len(entregas_filtradas)}")

    if not entregas_filtradas:
        print("No hay entregas que coincidan con el filtro para esta actividad.")
        return

    nombre_carpeta_actividad = construir_slug_actividad(coursework_title, coursework_id)

    # Si la carpeta_base ya apunta exactamente a esta actividad
    # (caso descarga de una sola actividad), reutilízala tal cual.
    if os.path.basename(os.path.normpath(carpeta_base)) == nombre_carpeta_actividad:
        carpeta_actividad = carpeta_base
    else:
        carpeta_actividad = os.path.join(carpeta_base, nombre_carpeta_actividad)

    # Evita folders duplicados cuando se reprocesa la misma actividad.
    # Si la carpeta ya existe, se reconstruye limpia con la convención actual.
    if os.path.exists(carpeta_actividad):
        shutil.rmtree(carpeta_actividad)

    asegurar_directorio(carpeta_actividad)

    for submission in entregas_filtradas:
        imprimir_resumen_entrega(submission)

        user_id = submission.get("userId", "sin_userId")

        if user_id not in perfiles_cache:
            perfil = {
                "correo": "",
                "nombre": "",
                "apellido": "",
            }

            if profile_scope_disponible:
                try:
                    perfil = obtener_perfil_usuario(
                        classroom_service=classroom_service,
                        user_id=user_id,
                        profile_scope_disponible=True,
                    )
                except HttpError as err:
                    print(
                        f"  ⚠️ No se pudo leer userProfiles para userId={user_id}. "
                        f"Se usará fallback. Detalle: {err}"
                    )
                    perfil = extraer_datos_usuario_desde_historial(submission)

            else:
                perfil = extraer_datos_usuario_desde_historial(submission)

            perfiles_cache[user_id] = perfil

        perfil = perfiles_cache[user_id]

        nombre_carpeta_entrega = construir_nombre_carpeta_entrega(
            submission=submission,
            perfil=perfil,
        )

        carpeta_entrega = os.path.join(carpeta_actividad, nombre_carpeta_entrega)
        asegurar_directorio(carpeta_entrega)

        rutas_descargadas: list[str] = []

        if se_puede_descargar_entrega(submission):
            rutas_descargadas = descargar_adjuntos_entrega(
                submission=submission,
                drive_service=drive_service,
                carpeta_entrega=carpeta_entrega,
            )

            if rutas_descargadas:
                estadisticas["archivos_descargados"] += len(rutas_descargadas)

            evaluacion = evaluar_entrega_automatica(
                submission=submission,
                rutas_descargadas=rutas_descargadas,
                coursework=coursework,
                autograding_config=AUTOGRADING_CONFIG,
            )
        else:
            print("  adjuntos: no aplica, alumno sin entrega enviada")
            evaluacion = {
                "auto_grade": 0,
                "feedback": "Alumno asignado pero sin entrega enviada. No se descargaron archivos ni se evaluó contenido.",
                "penalty_late": 0,
                "days_late": 0,
                "has_attachment": "false",
                "manual_review": "false",
                "readable_content": "false",
                "delivery_valid_score": 0,
                "evidence_score": 0,
                "readable_content_score": 0,
                "minimum_sufficiency_score": 0,
                "keyword_hits": "",
                "content_score": 0,
                "files_read_for_content": 0,
                "detected_words": 0,
            }

        print(f"  auto_grade: {evaluacion['auto_grade']}")
        print(f"  content_score: {evaluacion['content_score']}")
        print(f"  readable_content: {evaluacion['readable_content']}")
        print(f"  manual_review: {evaluacion['manual_review']}")
        print(f"  penalty_late: {evaluacion['penalty_late']}")

        filas_csv.append(
            {
                "curso": course_name,
                "actividad": coursework_title,
                "due_date": obtener_due_date_texto(coursework),
                "due_time": obtener_due_time_texto(coursework),
                "correo": perfil.get("correo", ""),
                "nombre": perfil.get("nombre", ""),
                "apellido": perfil.get("apellido", ""),
                "estado_entrega": estado_legible_entrega(submission),
                "late": str(bool(submission.get("late", False))).lower(),
                "assigned_grade": "" if submission.get("assignedGrade") is None else str(submission.get("assignedGrade")),
                "draft_grade": "" if submission.get("draftGrade") is None else str(submission.get("draftGrade")),
                "attached": str(tiene_adjuntos(submission)).lower(),
                "has_attachment": str(evaluacion["has_attachment"]).lower(),
                "manual_review": str(evaluacion["manual_review"]).lower(),
                "readable_content": str(evaluacion["readable_content"]).lower(),
                "days_late": str(evaluacion["days_late"]),
                "penalty_late": str(evaluacion["penalty_late"]),
                "delivery_valid_score": str(evaluacion["delivery_valid_score"]),
                "evidence_score": str(evaluacion["evidence_score"]),
                "readable_content_score": str(evaluacion["readable_content_score"]),
                "minimum_sufficiency_score": str(evaluacion["minimum_sufficiency_score"]),
                "keyword_hits": evaluacion["keyword_hits"],
                "content_score": str(evaluacion["content_score"]),
                "auto_grade": str(evaluacion["auto_grade"]),
                "feedback": evaluacion["feedback"],
            }
        )

        print("-" * 70)


# ==========================================================
# Flujo principal
# ==========================================================

AUTOGRADING_CONFIG = cargar_config_autograding()

def main() -> None:
    """
    Flujo principal:
    1. autentica
    2. elige curso
    3. elige alcance
    4. filtra actividades
    5. elige filtro de entregas
    6. elige formato de salida
    7. descarga
    8. genera CSV y zip
    """
    settings = get_settings()
    ensure_directories(settings)

    creds = get_credentials(
        credentials_path=settings.credentials_path,
        token_path=settings.token_path,
    )

    print("Autenticación correcta.")
    print(f"Token válido: {creds.valid}")

    try:
        classroom_service = build("classroom", "v1", credentials=creds)
        drive_service = build("drive", "v3", credentials=creds)

        profile_scope_disponible = detectar_scope_perfil(classroom_service)
        if profile_scope_disponible:
            print("✅ Scope de perfiles disponible.")
        else:
            print(
                "⚠️ El token no tiene scope para leer perfiles de alumnos. "
                "Se continuará con fallback y el CSV puede traer nombre/correo vacíos."
            )

        courses = obtener_todos_los_cursos(classroom_service)

        if not courses:
            print("No se encontraron cursos.")
            return

        selected_course = seleccionar_opcion(courses, "un curso")
        course_id = selected_course["id"]
        course_name = limpiar_nombre_archivo(selected_course.get("name", f"curso_{course_id}"))
        course_slug = construir_slug_curso(
            selected_course.get("name", f"curso_{course_id}"),
            course_id,
        )
        course_display = selected_course.get("display_name", selected_course.get("name", "Curso"))

        print(f"\n✅ Curso seleccionado: {course_display}")

        alcance_descarga = seleccionar_alcance_descarga()
        print(f"\n📚 Alcance seleccionado: {alcance_descarga}")

        filtro_actividades = seleccionar_filtro_actividades()
        print(f"\n Filtro de actividades: {filtro_actividades}")

        modo_descarga = seleccionar_modo_descarga()
        print(f"\n📥 Modo seleccionado: {describir_modo_descarga(modo_descarga)}")

        formato_salida = seleccionar_formato_salida()
        print(f"📦 Formato de salida: {formato_salida}")

        courseworks = obtener_todas_las_actividades(classroom_service, course_id)

        if not courseworks:
            print("No se encontraron actividades en este curso.")
            return

        courseworks_filtradas = filtrar_actividades(
            classroom_service=classroom_service,
            course_id=course_id,
            courseworks=courseworks,
            filtro=filtro_actividades,
        )

        if not courseworks_filtradas:
            print("No quedaron actividades después de aplicar el filtro.")
            return

        print(
            f"✅ Actividades encontradas: {len(courseworks)} | "
            f"después del filtro: {len(courseworks_filtradas)}"
        )

        carpeta_curso = os.path.normpath(
            os.path.join(settings.download_root, course_slug)
        )

        perfiles_cache: dict[str, dict[str, str]] = {}
        filas_csv: list[dict[str, str]] = []
        estadisticas = {
            "actividades_procesadas": 0,
            "entregas_totales": 0,
            "entregas_filtradas": 0,
            "archivos_descargados": 0,
        }

        if alcance_descarga == "single_coursework":
            selected_coursework = seleccionar_opcion(courseworks_filtradas, "una actividad")
            print(f"\n✅ Actividad seleccionada: {selected_coursework['display_name']}")

            coursework_title = selected_coursework.get("title", f"actividad_{selected_coursework['id']}")
            coursework_folder_name = construir_slug_actividad(
                coursework_title,
                selected_coursework["id"],
            )

            # Para una sola actividad también se reconstruye la carpeta del curso completa,
            # de modo que el ZIP conserve la estructura:
            # Curso_ID/Actividad_ID/alumno/
            carpeta_base = preparar_directorio_salida(
                carpeta_curso,
                limpiar_si_existe=True,
            )

            procesar_actividad(
                classroom_service=classroom_service,
                drive_service=drive_service,
                course=selected_course,
                coursework=selected_coursework,
                modo_descarga=modo_descarga,
                carpeta_base=carpeta_base,
                perfiles_cache=perfiles_cache,
                filas_csv=filas_csv,
                estadisticas=estadisticas,
                profile_scope_disponible=profile_scope_disponible,
            )

            nombre_csv = "resumen_entregas.csv"
            nombre_zip = course_slug

        elif alcance_descarga == "all_courseworks":
            # Para exportación completa sí se usa la carpeta del curso,
            # pero se limpia antes para que no arrastre histórico.
            carpeta_base = preparar_directorio_salida(
                carpeta_curso,
                limpiar_si_existe=True,
            )

            print(
                f"\n✅ Se procesarán todas las actividades filtradas del curso: "
                f"{len(courseworks_filtradas)}"
            )

            for idx, coursework in enumerate(courseworks_filtradas, start=1):
                print(f"\n[{idx}/{len(courseworks_filtradas)}] {coursework['display_name']}")

                procesar_actividad(
                    classroom_service=classroom_service,
                    drive_service=drive_service,
                    course=selected_course,
                    coursework=coursework,
                    modo_descarga=modo_descarga,
                    carpeta_base=carpeta_base,
                    perfiles_cache=perfiles_cache,
                    filas_csv=filas_csv,
                    estadisticas=estadisticas,
                    profile_scope_disponible=profile_scope_disponible,
                )

            nombre_csv = "resumen_todas_las_actividades.csv"
            nombre_zip = course_slug

        else:
            print("❌ Alcance de descarga no reconocido.")
            return

        csv_path = os.path.join(carpeta_base, nombre_csv)
        escribir_csv_resumen(csv_path, filas_csv)

        if formato_salida == "zip_and_folder":
            zip_base_name = os.path.join("downloads", nombre_zip)
            comprimir_carpeta_a_zip(carpeta_base, zip_base_name)

        print("\n" + "=" * 90)
        print("RESUMEN FINAL")
        print("=" * 90)
        print(f"Curso: {course_display}")
        print(f"Actividades procesadas: {estadisticas['actividades_procesadas']}")
        print(f"Entregas totales vistas: {estadisticas['entregas_totales']}")
        print(f"Entregas que cumplieron filtro: {estadisticas['entregas_filtradas']}")
        print(f"Archivos descargados: {estadisticas['archivos_descargados']}")
        print(f"Filas en CSV: {len(filas_csv)}")
        print("Nota: en modo 'all' ahora el CSV incluye también alumnos sin entregar.")
        print(f"Carpeta base: {carpeta_base}")
        if formato_salida == "zip_and_folder":
            print(f"ZIP: downloads/{nombre_zip}.zip")

        print("\n✅ Proceso terminado.")

    except HttpError as err:
        print(f"Error al consultar Classroom: {err}")


if __name__ == "__main__":
    main()
